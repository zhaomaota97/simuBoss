from __future__ import annotations

import copy
import json
import sys
from pathlib import Path

from pptx import Presentation


BASE_DIR = Path(__file__).resolve().parent
TEMPLATE = BASE_DIR / "caika-template.pptx"
CONTENT = BASE_DIR / "content.json"
OUTPUT = BASE_DIR / "财咖分析云售前SOP实战指南-生成版.pptx"

SECTION_COVER_INDICES = [2, 5, 13, 17, 21]
LAYOUTS = {
    "cover": 0,
    "toc": 1,
    "focus_2": 3,
    "levels_3": 4,
    "insight_4": 6,
    "duo_detail": 7,
    "numbered_3": 8,
    "process_4": 9,
    "pair_cards_2": 10,
    "columns_3": 11,
    "balanced_4": 12,
    "dual_story": 14,
    "stacked_3": 15,
    "quadrant_4": 16,
    "compare_2": 18,
    "advantage_4": 19,
    "service_3": 20,
    "followup_4": 22,
    "followup_2": 23,
    "followup_3": 24,
    "ending": 25,
}

PML_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"


def _configure_stdio() -> None:
    for stream_name in ("stdout", "stderr"):
        stream = getattr(sys, stream_name, None)
        if stream is None or not hasattr(stream, "reconfigure"):
            continue
        try:
            stream.reconfigure(errors="backslashreplace")
        except Exception:
            pass


def load_content(path: str | Path = CONTENT) -> dict:
    with Path(path).open("r", encoding="utf-8") as handle:
        return json.load(handle)


def get_shape(slide, shape_name: str):
    for shape in slide.shapes:
        if shape.name == shape_name and hasattr(shape, "text_frame"):
            return shape
    return None


def set_shape_text(slide, shape_name: str, text: str) -> None:
    shape = get_shape(slide, shape_name)
    if shape is None:
        return

    frame = shape.text_frame
    value = str(text or "")
    paragraphs = list(frame.paragraphs)
    if not paragraphs:
        frame.text = value
        return

    first_paragraph = paragraphs[0]
    if first_paragraph.runs:
        first_paragraph.runs[0].text = value
        for run in list(first_paragraph.runs[1:]):
            first_paragraph._p.remove(run._r)
    else:
        run = first_paragraph.add_run()
        run.text = value

    for paragraph in list(frame.paragraphs[1:]):
        paragraph._p.getparent().remove(paragraph._p)


def _remap_relationship_ids(element, rel_map: dict[str, str]) -> None:
    for node in element.iter():
        for attr_name, attr_value in list(node.attrib.items()):
            if attr_value in rel_map:
                node.set(attr_name, rel_map[attr_value])


def copy_slide(prs: Presentation, source_index: int):
    source = prs.slides[source_index]
    new_slide = prs.slides.add_slide(source.slide_layout)

    for shape in list(new_slide.shapes):
        sp = shape.element
        sp.getparent().remove(sp)

    rel_map = {}
    for rel_id, rel in source.part.rels.items():
        if "slideLayout" in rel.reltype or "notesSlide" in rel.reltype:
            continue
        if rel.is_external:
            new_rel_id = new_slide.part.rels._add_relationship(rel.reltype, rel.target_ref, is_external=True)
        else:
            new_rel_id = new_slide.part.rels._add_relationship(rel.reltype, rel.target_part, is_external=False)
        rel_map[rel_id] = new_rel_id

    source_c_sld = source.element.find(f"{{{PML_NS}}}cSld")
    target_c_sld = new_slide.element.find(f"{{{PML_NS}}}cSld")
    if source_c_sld is not None and target_c_sld is not None:
        source_bg = source_c_sld.find(f"{{{PML_NS}}}bg")
        target_bg = target_c_sld.find(f"{{{PML_NS}}}bg")
        if target_bg is not None:
            target_c_sld.remove(target_bg)
        if source_bg is not None:
            new_bg = copy.deepcopy(source_bg)
            _remap_relationship_ids(new_bg, rel_map)
            target_c_sld.insert(0, new_bg)

    for shape in source.shapes:
        new_el = copy.deepcopy(shape.element)
        _remap_relationship_ids(new_el, rel_map)
        new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")

    return new_slide


def delete_original_slides(prs: Presentation, original_count: int) -> None:
    slide_ids = prs.slides._sldIdLst  # noqa: SLF001
    for _ in range(original_count):
        element = slide_ids[0]
        slide_ids.remove(element)


def _merge_text(*parts: str) -> str:
    cleaned = [str(part).strip() for part in parts if str(part or "").strip()]
    return "\n".join(cleaned)


def _normalize_section_id(value: str | None, fallback_index: int) -> str:
    raw = str(value or "").strip()
    digits = "".join(ch for ch in raw if ch.isdigit())
    if digits:
        return digits.zfill(2)
    return f"{fallback_index + 1:02d}"


def _section_titles(content_data: dict) -> list[str]:
    sections = list(content_data.get("sections") or [])
    return [str(section.get("title") or "").strip() for section in sections[:5]]


def _items(page: dict, key: str = "items") -> list[dict]:
    return list(page.get(key) or [])


def _set_triples(slide, slots: list[tuple[str, str, str]], items: list[dict]) -> None:
    for index, (a_name, b_name, c_name) in enumerate(slots):
        item = items[index] if index < len(items) else {}
        set_shape_text(slide, a_name, item.get("label", "") if item else "")
        set_shape_text(slide, b_name, item.get("title", ""))
        set_shape_text(slide, c_name, item.get("text", ""))


def _set_pairs(slide, slots: list[tuple[str, str]], items: list[dict], use_note: bool = False) -> None:
    for index, (title_name, body_name) in enumerate(slots):
        item = items[index] if index < len(items) else {}
        body = _merge_text(item.get("text", ""), item.get("note", "")) if use_note else item.get("text", "")
        set_shape_text(slide, title_name, item.get("title", ""))
        set_shape_text(slide, body_name, body)


def build_cover(prs: Presentation, data: dict):
    slide = copy_slide(prs, LAYOUTS["cover"])
    set_shape_text(slide, "Text 7", data.get("title", ""))
    set_shape_text(slide, "Text 1", data.get("subtitle", ""))
    set_shape_text(slide, "Text 2", data.get("date", ""))
    return slide


def build_toc(prs: Presentation, section_titles: list[str]):
    slide = copy_slide(prs, LAYOUTS["toc"])
    slots = [
        ("Text 8", "Text 9"),
        ("Text 10", "Text 11"),
        ("Text 12", "Text 13"),
        ("Text 14", "Text 15"),
        ("Text 16", "Text 17"),
    ]
    for index, (number_name, title_name) in enumerate(slots):
        title = section_titles[index] if index < len(section_titles) else ""
        set_shape_text(slide, number_name, f"{index + 1:02d}." if title else "")
        set_shape_text(slide, title_name, title)
    return slide


def build_section_cover(prs: Presentation, section: dict, section_index: int):
    template_index = SECTION_COVER_INDICES[min(section_index, len(SECTION_COVER_INDICES) - 1)]
    slide = copy_slide(prs, template_index)
    set_shape_text(slide, "Text 0", _normalize_section_id(section.get("id"), section_index))
    set_shape_text(slide, "Text 1", section.get("title", ""))
    return slide


def build_focus_2(prs: Presentation, page: dict):
    slide = copy_slide(prs, LAYOUTS["focus_2"])
    set_shape_text(slide, "Text 0", page.get("title", ""))
    _set_pairs(slide, [("Text 5", "Text 6"), ("Text 7", "Text 8")], _items(page), False)
    return slide


def build_levels_3(prs: Presentation, page: dict):
    slide = copy_slide(prs, LAYOUTS["levels_3"])
    set_shape_text(slide, "Text 0", page.get("title", ""))
    _set_triples(
        slide,
        [("Text 1", "Text 2", "Text 3"), ("Text 4", "Text 5", "Text 6"), ("Text 7", "Text 8", "Text 9")],
        _items(page),
    )
    return slide


def build_insight_4(prs: Presentation, page: dict):
    slide = copy_slide(prs, LAYOUTS["insight_4"])
    set_shape_text(slide, "Text 0", page.get("title", ""))
    _set_pairs(
        slide,
        [("Text 5", "Text 6"), ("Text 7", "Text 8"), ("Text 9", "Text 10"), ("Text 11", "Text 12")],
        _items(page),
        False,
    )
    return slide


def build_duo_detail(prs: Presentation, page: dict):
    slide = copy_slide(prs, LAYOUTS["duo_detail"])
    set_shape_text(slide, "Text 0", page.get("title", ""))
    _set_pairs(slide, [("Text 5", "Text 6"), ("Text 7", "Text 8")], _items(page), False)
    return slide


def build_numbered_3(prs: Presentation, page: dict):
    slide = copy_slide(prs, LAYOUTS["numbered_3"])
    set_shape_text(slide, "Text 0", page.get("title", ""))
    _set_triples(
        slide,
        [("Text 1", "Text 2", "Text 3"), ("Text 4", "Text 5", "Text 6"), ("Text 7", "Text 8", "Text 9")],
        _items(page),
    )
    return slide


def build_process_4(prs: Presentation, page: dict):
    slide = copy_slide(prs, LAYOUTS["process_4"])
    set_shape_text(slide, "Text 0", page.get("title", ""))
    _set_triples(
        slide,
        [("Text 4", "Text 5", "Text 6"), ("Text 7", "Text 8", "Text 9"), ("Text 10", "Text 11", "Text 12"), ("Text 13", "Text 14", "Text 15")],
        _items(page, "steps"),
    )
    return slide


def build_pair_cards_2(prs: Presentation, page: dict):
    slide = copy_slide(prs, LAYOUTS["pair_cards_2"])
    set_shape_text(slide, "Text 0", page.get("title", ""))
    _set_triples(slide, [("Text 7", "Text 8", "Text 9"), ("Text 10", "Text 11", "Text 12")], _items(page))
    return slide


def build_columns_3(prs: Presentation, page: dict):
    slide = copy_slide(prs, LAYOUTS["columns_3"])
    set_shape_text(slide, "Text 0", page.get("title", ""))
    _set_pairs(slide, [("Text 7", "Text 8"), ("Text 9", "Text 10"), ("Text 11", "Text 12")], _items(page), False)
    return slide


def build_balanced_4(prs: Presentation, page: dict):
    slide = copy_slide(prs, LAYOUTS["balanced_4"])
    set_shape_text(slide, "Text 0", page.get("title", ""))
    _set_pairs(slide, [("Text 5", "Text 6"), ("Text 7", "Text 8"), ("Text 9", "Text 10"), ("Text 11", "Text 12")], _items(page), False)
    return slide


def build_dual_story(prs: Presentation, page: dict):
    slide = copy_slide(prs, LAYOUTS["dual_story"])
    set_shape_text(slide, "Text 0", page.get("title", ""))
    _set_pairs(slide, [("Text 5", "Text 6"), ("Text 7", "Text 8")], _items(page), False)
    return slide


def build_stacked_3(prs: Presentation, page: dict):
    slide = copy_slide(prs, LAYOUTS["stacked_3"])
    set_shape_text(slide, "Text 0", page.get("title", ""))
    _set_triples(
        slide,
        [("Text 1", "Text 2", "Text 3"), ("Text 4", "Text 5", "Text 6"), ("Text 7", "Text 8", "Text 9")],
        _items(page),
    )
    return slide


def build_quadrant_4(prs: Presentation, page: dict):
    slide = copy_slide(prs, LAYOUTS["quadrant_4"])
    set_shape_text(slide, "Text 0", page.get("title", ""))
    _set_pairs(slide, [("Text 2", "Text 3"), ("Text 4", "Text 5"), ("Text 6", "Text 7"), ("Text 8", "Text 9")], _items(page), False)
    return slide


def build_compare_2(prs: Presentation, page: dict):
    slide = copy_slide(prs, LAYOUTS["compare_2"])
    set_shape_text(slide, "Text 0", page.get("title", ""))
    _set_triples(slide, [("Text 5", "Text 6", "Text 7"), ("Text 8", "Text 9", "Text 10")], _items(page))
    return slide


def build_advantage_4(prs: Presentation, page: dict):
    slide = copy_slide(prs, LAYOUTS["advantage_4"])
    set_shape_text(slide, "Text 0", page.get("title", ""))
    _set_triples(
        slide,
        [("Text 8", "Text 9", "Text 10"), ("Text 11", "Text 12", "Text 13"), ("Text 14", "Text 15", "Text 16"), ("Text 17", "Text 18", "Text 19")],
        _items(page),
    )
    return slide


def build_service_3(prs: Presentation, page: dict):
    slide = copy_slide(prs, LAYOUTS["service_3"])
    set_shape_text(slide, "Text 0", page.get("title", ""))
    _set_pairs(slide, [("Text 5", "Text 6"), ("Text 7", "Text 8"), ("Text 9", "Text 10")], _items(page), False)
    return slide


def build_followup_4(prs: Presentation, page: dict):
    slide = copy_slide(prs, LAYOUTS["followup_4"])
    set_shape_text(slide, "Text 0", page.get("title", ""))
    _set_pairs(slide, [("Text 1", "Text 2"), ("Text 3", "Text 4"), ("Text 5", "Text 6"), ("Text 7", "Text 8")], _items(page), False)
    return slide


def build_followup_2(prs: Presentation, page: dict):
    slide = copy_slide(prs, LAYOUTS["followup_2"])
    set_shape_text(slide, "Text 0", page.get("title", ""))
    _set_triples(slide, [("Text 6", "Text 7", "Text 8"), ("Text 9", "Text 10", "Text 11")], _items(page))
    return slide


def build_followup_3(prs: Presentation, page: dict):
    slide = copy_slide(prs, LAYOUTS["followup_3"])
    set_shape_text(slide, "Text 0", page.get("title", ""))
    _set_pairs(slide, [("Text 3", "Text 4"), ("Text 5", "Text 6"), ("Text 7", "Text 8")], _items(page), False)
    return slide


def build_ending(prs: Presentation, data: dict):
    slide = copy_slide(prs, LAYOUTS["ending"])
    set_shape_text(slide, "Text 4", data.get("thanks", "THANK YOU FOR READING！"))
    set_shape_text(slide, "Text 2", data.get("message", "财务小变化、企业大未来"))
    set_shape_text(slide, "Text 6", data.get("name", ""))
    set_shape_text(slide, "Text 7", data.get("date", ""))
    return slide


def generate(template_path: str | Path, content_data: dict, output_path: str | Path) -> Path:
    prs = Presentation(str(template_path))
    original_count = len(prs.slides)

    build_cover(prs, content_data.get("cover") or {})
    build_toc(prs, _section_titles(content_data))

    builders = {
        "focus_2": build_focus_2,
        "levels_3": build_levels_3,
        "insight_4": build_insight_4,
        "duo_detail": build_duo_detail,
        "numbered_3": build_numbered_3,
        "process_4": build_process_4,
        "pair_cards_2": build_pair_cards_2,
        "columns_3": build_columns_3,
        "balanced_4": build_balanced_4,
        "dual_story": build_dual_story,
        "stacked_3": build_stacked_3,
        "quadrant_4": build_quadrant_4,
        "compare_2": build_compare_2,
        "advantage_4": build_advantage_4,
        "service_3": build_service_3,
        "followup_4": build_followup_4,
        "followup_2": build_followup_2,
        "followup_3": build_followup_3,
    }

    sections = list(content_data.get("sections") or [])
    for section_index, section in enumerate(sections):
        build_section_cover(prs, section, section_index)
        for page in section.get("pages") or []:
            layout = str(page.get("layout") or "").strip().lower()
            builder = builders.get(layout)
            if builder is None:
                print(f"[WARN] Unknown page layout: {layout}")
                continue
            builder(prs, page)

    build_ending(prs, content_data.get("ending") or {})
    delete_original_slides(prs, original_count)

    final_path = Path(output_path)
    final_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(final_path))
    print(f"[OK] PPT generated: {final_path.name}")
    return final_path


def generate_from_file(
    template_path: str | Path = TEMPLATE,
    content_path: str | Path = CONTENT,
    output_path: str | Path = OUTPUT,
) -> Path:
    return generate(template_path, load_content(content_path), output_path)


if __name__ == "__main__":
    _configure_stdio()
    generate_from_file()
